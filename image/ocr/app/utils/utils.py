import base64
import io
import os
from PIL import Image, ImageSequence
from pdf2image import convert_from_bytes
import numpy as np
import gc
import shutil
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from bs4 import BeautifulSoup as BS
import copy
from fastapi import FastAPI
from detectaicore import Job, image_file_names
import zipfile
from pathlib import Path
import logging
import pymupdf
from rapidocr_onnxruntime import RapidOCR, VisRes
from .ocr_c import pil2PIX32
import subprocess

from PIL import ImageFile

ImageFile.LOAD_TRUNCATED_IMAGES = True

ROOT_DIR = os.path.dirname(os.path.abspath(__file__))


def create_engine_rapidocr():
    """
    Create RapidOCR Engine
    return:
    RapidOCR Engine
    """
    # create RapidOCR Engine
    engine = RapidOCR(det_model_type="DB", rec_model_type="CRNN")
    return engine


def ocr_to_text(ocr_output):
    """
    # Sort the OCR output based on the y-coordinate of the bounding box
    each element of the list is another list wihich contain 3 elements

        element 0 bounding box coordinates

        element 1 Text extracted

        element 2 score

        Example:

        [[[529.0, 149.0], [577.0, 149.0], [577.0, 174.0], [529.0, 174.0]],
        'R&D',
        0.9749892751375834]
    Args:
        ocr_output : list of lists, output from Rapid OCR
    """
    # Sort the OCR output based on the y-coordinate of the bounding box
    sorted_ocr = sorted(ocr_output, key=lambda x: x[0][0][1])

    # Initialize variables
    lines = []
    current_line = []
    current_y = sorted_ocr[0][0][0][1]

    # Iterate through the sorted OCR output
    for item in sorted_ocr:
        bbox, text, score = item
        bbox_y = bbox[0][1]

        # Check if the bounding box is on the same line as the current line
        if abs(bbox_y - current_y) <= 10:
            current_line.append(text + " ")
        else:
            lines.append(" ".join(current_line))
            current_line = [text + " "]
            current_y = bbox_y

    # Append the last line
    if current_line:
        lines.append(" ".join(current_line))

    # Join all lines into a single text output
    text_output = "\n".join(lines)
    return text_output


def create_directory_if_not_exists(directory_path):
    """
    Creates a directory if it does not already exist.

    Args:
      directory_path: The path to the directory to create.

    Raises:
      OSError: If an error occurs during directory creation.
    """
    try:
        os.makedirs(directory_path, exist_ok=True)  # Create the directory

    except OSError as e:
        logging.error(f"Error creating directory: {e}")


def remove_file_if_exists(file_path):
    """
    Removes a file if it exists.

    Args:
      file_path: The path to the file to be removed.

    Returns:
      True if the file was successfully removed, False otherwise.
    """
    try:
        if os.path.exists(file_path):
            os.remove(file_path)

        else:

            return False
    except OSError as e:
        logging.error(f"Error removing file: {e}")
        return False


def ocr_image(im, app):
    """
    Extract text from an Image using Leptonica and tessetact
    params:
    im: PIL image
    app: FastAPI endpoint. Contains global objects with tesseract and Leptonica C objects
    return:
    text extracted
    """
    if os.getenv("USE_RAPID_OCR") == "YES":
        engine = app.engine
        # Perform OCR
        directory_to_create = os.path.join(
            os.path.dirname(os.path.abspath(__file__)), "temp"
        )
        create_directory_if_not_exists(directory_to_create)
        tempfile = os.path.join(directory_to_create, "temp.png")
        im = im.convert("RGB")
        im.save(tempfile)
        result, _ = engine(tempfile)
        # Extract text while preserving spaces and returns
        extracted_text = ocr_to_text(result)  # "\n".join([line[1] for line in result])
        remove_file_if_exists(tempfile)
        del im
        gc.collect()
        return extracted_text
    else:

        pix = pil2PIX32(im, app.leptonica, app.ffi)

        # Get information about DPI
        x_dpi = app.ffi.new("int *")
        y_dpi = app.ffi.new("int *")
        app.leptonica.pixGetResolution(pix, x_dpi, y_dpi)
        app.tesseract.TessBaseAPISetImage2(app.api, pix)
        app.tesseract.TessBaseAPIRecognize(app.api, app.ffi.NULL)
        # Print whole recognized text
        utf8_text = app.ffi.string(
            app.tesseract.TessBaseAPIGetUTF8Text(app.api)
        ).decode("utf-8")
        del pix
        del x_dpi
        del y_dpi
        gc.collect()
        return utf8_text


def workout_git(im):
    """
    Create RGBA image from a GIF image
    params:
    im: PIL image
    """
    mypalette = im.getpalette()
    im.putpalette(mypalette)
    new_im = Image.new("RGBA", im.size)
    new_im.paste(im)
    im = copy.deepcopy(new_im)
    del new_im
    gc.collect()
    return im


def zoom_png(im, zoom_x=2.0, zoom_y=2.0):
    """
    Zoom in a PNG image
    params:
    im: PIL image
    """
    logging.info(f"Zooming Image")
    doc = pymupdf.open(stream=io.BytesIO(im))
    page = doc[0]
    # get pixel map
    pix = page.get_pixmap()
    # Zoom in the image
    zoom_x = zoom_x  # horizontal zoom
    zoom_y = zoom_y  # vertical zoom
    mat = pymupdf.Matrix(zoom_x, zoom_y)  # zoom factor 2 in each dimension
    pix = page.get_pixmap(matrix=mat)
    im = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
    im = im.convert("RGB")
    del doc
    del page
    del pix
    gc.collect()
    return im


def extract_from_single_image(data, app):
    """
    Process imagenes and extract text via ocr
    params:
    data: list of documents from detect AI
    app global variables from FastAPI. contains tessetact and leptonica objects
    return:
    text extracted and endpoint
    """
    im_b64 = data.get("source").get("content")
    ext = data.get("source").get("file_type")

    img_bytes = base64.b64decode(im_b64.encode("utf-8"))

    logging.info(f"Extract Image type: {ext}, Length: {len(img_bytes)}")
    im = zoom_png(img_bytes)
    # im = Image.open(io.BytesIO(img_bytes))
    # if ext == "gif":
    #     # get RGBA IMage
    #     im = Image.open(io.BytesIO(img_bytes))
    #     im = workout_git(im)
    # else:
    #     im = zoom_png(img_bytes)

    logging.info(f"Extract Image type: {ext}, Size: {im.size}")
    # extract text
    utf8_text = ocr_image(im, app)
    logging.info(f"Extract Image type: {ext}, Length text: {len(utf8_text)}")
    data["source"]["content"] = utf8_text
    del im
    del img_bytes
    del im_b64
    gc.collect()
    return utf8_text, data


def extract_text_from_pdf(data, app):
    """
    extract and Process images in a pdf file and extract text via ocr
    params:
    data: list of documents from detect AI
    app global variables from FastAPI. contains tessetact and leptonica objects
    return:
    text extracted and endpoint
    """

    text = {}
    im_b64 = data.get("source").get("content")

    img_bytes = base64.b64decode(im_b64.encode("utf-8"))
    logging.info(f"Extract Image type: PDF, Length: {len(img_bytes)}")
    # Open Image
    pdf_file = convert_from_bytes(img_bytes, dpi=300)
    for i, page in enumerate(pdf_file):
        page_arr = np.asarray(page)
        im = Image.fromarray(page_arr)
        # Extract text from Image
        utf8_text = ocr_image(im, app)
        text[i] = utf8_text
        del im
        del page_arr
        del utf8_text
    # Final text in single file with a mark
    txt = ""
    mark = "\n### NEW PAGE ###\n"
    txt = mark.join(text[key] for key in text.keys()) + mark

    logging.info(
        f"Extract Image type: PDF, number pages: {len(text.keys())}, Length text: {len(txt[: -len(mark)])}"
    )
    del pdf_file
    del img_bytes
    del im_b64
    del text
    gc.collect()
    return txt[: -len(mark)]


def extract_text_from_tiff(data, app):
    """
    extract and Process images in a tiff/tif file
    params:
    data: list of documents from detect AI
    app global variables from FastAPI. contains tessetact and leptonica objects
    return:
    text extracted and endpoint
    """

    text = {}
    im_b64 = data.get("source").get("content")

    img_bytes = base64.b64decode(im_b64.encode("utf-8"))
    logging.info(f"Extract Image type: TIF, Length: {len(img_bytes)}")
    # Open Image
    imfile = Image.open(io.BytesIO(img_bytes))

    for i, page in enumerate(ImageSequence.Iterator(imfile)):
        im = page.convert("RGB")
        # Extract text from Image
        utf8_text = ocr_image(im, app)
        text[i] = utf8_text
        del im
        del utf8_text
        gc.collect()
    # Final text in single file with a mark
    txt = ""
    mark = "\n### NEW PAGE ###\n"
    txt = mark.join(text[key] for key in text.keys()) + mark

    logging.info(
        f"Extract Image type: TIFF/TIF, number pages: {len(text.keys())}, Length text: {len(txt[: -len(mark)])}"
    )
    del imfile
    del img_bytes
    del im_b64
    del text
    gc.collect()
    return txt[: -len(mark)]


def extract_images_from_docx(docx_path, output_folder):
    """
    # Ensure the output folder exists
    Args:
        docx_path: Path to the input .docx file.
        output_folder: Path to the output folder where the images will be extracted.
    """
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    # Open the DOCX file as a zip file
    with zipfile.ZipFile(docx_path, "r") as docx_zip:
        # Iterate through the files in the zip
        for file in docx_zip.namelist():
            # Check if the file is an image
            if file.startswith("word/media/"):
                # Extract the image to the output folder
                docx_zip.extract(file, output_folder)
                logging.info(f"Extracted {file} to {output_folder}")
    return


def list_jpg_files(folder_path):
    """
    list folder and subfolder for jpg files
    Args:
        folder_path: Path to the folder to search for JPG files.
    """
    jpg_files = []
    # Walk through the folder and its subfolders
    for root, _, files in os.walk(folder_path):
        for file in files:
            # Check if the file is a JPG
            if file.lower().endswith(".jpg") or file.lower().endswith(".jpeg"):
                jpg_files.append(os.path.join(root, file))
    return jpg_files


# def convert_doc_to_docx(doc_path, docx_path):
#     """
#     # Convert DOC to DOCX using unoconv
#     Args:
#         doc_path: Path to the input .doc file.
#         docx_path: Path to the output
#     """
#     command = ["unoconv", "-f", "docx", "-o", docx_path, doc_path]
#     process = subprocess.Popen(command)
#     try:
#         process.wait(timeout=120)
#         logging.info(f"Converted {doc_path} to {docx_path} susccessfully")
#         process.kill()
#         logging.warning(f"kill process {doc_path} to {docx_path} susccessfully")
#     except subprocess.TimeoutExpired:
#         # Kill the process if it exceeds the timeout
#         process.kill()
#         logging.error(f"Error converting {doc_path} to {docx_path}: Timeout")

#     return


def convert_doc_to_docx(doc_path, docx_path, directory):
    """
    # Convert DOC to DOCX using LibreOffice command line
    Args:
        doc_path: Path to the input .doc file.
        docx_path: Path to the output .docx file.
    """
    command = [
        "libreoffice",
        "--headless",
        "--convert-to",
        "docx",
        "--outdir",
        os.path.dirname(docx_path),
        doc_path,
    ]
    process = subprocess.Popen(command)
    try:
        process.wait(timeout=20)
        logging.info(
            f"Converted {doc_path} to {docx_path} susccessfully. convert_doc_to_docx Function"
        )
        process.kill()
    except subprocess.TimeoutExpired:
        # Kill the process if it exceeds the timeout
        process.kill()
        shutil.rmtree(directory, ignore_errors=True)
        logging.error(f"Error converting {doc_path} to {docx_path}: Timeout")
    return


def extract_text_from_doc(data, app):
    """
    extract and Process images in a doc/docx file and extract text via ocr
    params:
    data: list of documents from detect AI
    app global variables from FastAPI. contains tessetact and leptonica objects
    return:
    text extracted and endpoint
    """
    # temporal directories
    directory = os.path.join(ROOT_DIR, "tmp")
    if not os.path.exists(directory):
        os.makedirs(directory)

    images = os.path.join(ROOT_DIR, "tmp", "images")
    if not os.path.exists(images):
        os.makedirs(images)

    text = {}
    im_b64 = data.get("source").get("content")
    ext = data.get("source").get("file_type")
    if ext == "docx":
        # write docx to temporal folder
        temp_docx = os.path.join(directory, "temp.docx")

        with open(temp_docx, "wb") as f:
            f.write(base64.b64decode(im_b64))
    else:
        temp_doc = os.path.join(directory, "temp.doc")
        with open(temp_doc, "wb") as f:
            f.write(base64.b64decode(im_b64))
        temp_docx = os.path.join(directory, "temp.docx")
        # convert DOC to DOCX
        convert_doc_to_docx(temp_doc, temp_docx, directory)

    logging.info(
        f"Extract Image type: DOC/DOCX, Length: {len(base64.b64decode(im_b64))}"
    )
    # _ = docx2txt.process(temp_docx, images)
    extract_images_from_docx(temp_docx, images)
    files = list_jpg_files(images)

    for i, file in zip(range(len(files)), files):
        logging.info(f"Extracting text from {file}")
        im = Image.open(file)
        # Extract text from Image
        utf8_text = ocr_image(im, app)
        text[i] = utf8_text
        del im
        del utf8_text
        gc.collect()
    # Final text in single file with a mark
    txt = ""
    mark = "\n### NEW PAGE ###\n"
    txt = mark.join(text[key] for key in text.keys()) + mark

    # delete temp folder
    shutil.rmtree(directory, ignore_errors=True)
    logging.info(
        f"Extract Image type: DOC/DOCX, number pages: {len(text.keys())}, Length text: {len(txt[: -len(mark)])}"
    )

    del files
    del im_b64
    del text
    gc.collect()
    return txt[: -len(mark)]


def convert_ppt_to_pptx_with_soffice(ppt_path, output_dir):
    """
    Converts a .ppt file to .pptx format using the soffice command-line tool.

    Args:
      ppt_path: Path to the input .ppt file.
      output_dir: Path to the output directory.

    Returns:
      True if the conversion was successful, False otherwise.
    """
    command = [
        "soffice",
        "--headless",
        "--convert-to",
        "pptx",
        "--outdir",
        output_dir,
        ppt_path,
    ]
    process = subprocess.Popen(command)
    try:
        process.wait(timeout=20)
        logging.info(
            f"Converted {ppt_path} to .pptx in {output_dir}. convert_ppt_to_pptx_with_soffice Function"
        )
        process.kill()

    except subprocess.TimeoutExpired:
        # Kill the process if it exceeds the timeout
        process.kill()
        logging.error(f"Error converting {ppt_path} to .pptx: Timeout")
    return


def extract_text_from_ppt(data, app):
    """
    extract and Process images in a ppt/pptx file and extract text via ocr
    params:
    data: list of documents from detect AI
    app global variables from FastAPI. contains tessetact and leptonica objects
    return:
    text extracted and endpoint
    """

    def write_image(shape, images, n):
        image = shape.image
        # ---get image "file" contents---
        image_bytes = image.blob
        # ---make up a name for the file, e.g. 'image.jpg'---
        image_filename = os.path.join(images, "ppt_image{:03d}.{}".format(n, image.ext))
        n += 1

        with open(image_filename, "wb") as f:
            f.write(image_bytes)
        return n

    def visitor(shape, images, n):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for s in shape.shapes:
                n = visitor(s, images, n)
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            n = write_image(shape, images, n)
        return n

    def iter_picture_shapes(prs, n):
        for slide in prs.slides:
            for shape in slide.shapes:
                n = visitor(shape, images, n)

    # page number
    n = 0
    # temporal directories
    directory = os.path.join(ROOT_DIR, "tmp")
    if not os.path.exists(directory):
        os.makedirs(directory)

    images = os.path.join(ROOT_DIR, "tmp", "images")
    if not os.path.exists(images):
        os.makedirs(images)

    text = {}
    im_b64 = data.get("source").get("content")
    ext = data.get("source").get("file_type")
    # if ppt convert to pptx with soffice
    if ext == "pptx":
        # write docx to temporal folder
        temp_pptx = os.path.join(directory, "temp.pptx")

        with open(temp_pptx, "wb") as f:
            f.write(base64.b64decode(im_b64))
    elif ext == "ppt":
        # write docx to temporal folder
        temp_ppt = os.path.join(directory, "temp.ppt")

        with open(temp_ppt, "wb") as f:
            f.write(base64.b64decode(im_b64))
        temp_pptx = os.path.join(directory, "temp.pptx")
        # convert PPT to PPTX
        convert_ppt_to_pptx_with_soffice(temp_ppt, directory)

    logging.info(
        f"Extract Image type: PPT/PPTX, Length: {len(base64.b64decode(im_b64))}"
    )
    # Open PPT and extract Images
    iter_picture_shapes(Presentation(temp_pptx), n)
    # list all images extracted in images
    files = os.listdir(images)

    for i, file in zip(range(len(files)), files):
        im = Image.open(os.path.join(images, file))
        # Extract text from Image
        utf8_text = ocr_image(im, app)
        text[i] = utf8_text
        del im
        del utf8_text
        gc.collect()
    # Final text in single file with a mark
    txt = ""
    mark = "\n### NEW PAGE ###\n"
    txt = mark.join(text[key] for key in text.keys()) + mark

    # delete temp folder
    shutil.rmtree(directory, ignore_errors=True)

    logging.info(
        f"Extract Image type: PPT/PPTX, number pages: {len(text.keys())}, Length text: {len(txt[: -len(mark)])}"
    )
    del files
    del im_b64
    del text
    gc.collect()
    return txt[: -len(mark)]


def extract_text_from_html(data, app):
    """
    extract and Process images in a html file with images embedded in base64 in src and extract text via ocr
    params:
    data: list of documents from detect AI
    app global variables from FastAPI. contains tessetact and leptonica objects
    return:
    text extracted and endpoint
    """

    # page number
    n = 0
    # temporal directories
    directory = os.path.join(ROOT_DIR, "tmp")
    if not os.path.exists(directory):
        os.makedirs(directory)

    images = os.path.join(ROOT_DIR, "tmp", "images")
    if not os.path.exists(images):
        os.makedirs(images)

    text = {}
    im_b64 = data.get("source").get("content")
    # write docx to temporal folder
    html_temp = os.path.join(directory, "temp.html")

    with open(html_temp, "wb") as f:
        f.write(base64.b64decode(im_b64))

    logging.info(f"Extract Image type: HTML, Length: {len(base64.b64decode(im_b64))}")
    # Open html and extract Images
    with open(html_temp) as html_wr:
        html_data = html_wr.read()

    soup = BS(html_data, features="lxml")

    images_arr = soup.find_all("img")

    for ind, i in zip(images_arr, range(len(images_arr))):
        image_data_base64 = ind["src"].split(",")[1]
        decoded_img_data = base64.b64decode(image_data_base64)
        with open(os.path.join(images, f"site_{i}.png"), "wb+") as img_wr:
            img_wr.write(decoded_img_data)
        del decoded_img_data
        gc.collect()

    logging.info(f"Extract Image type: HTML, number images: {len(images_arr)}")

    # list all images extracted in images
    files = os.listdir(images)

    for i, file in zip(range(len(files)), files):
        im = Image.open(os.path.join(images, file))
        # Extract text from Image
        utf8_text = ocr_image(im, app)
        text[i] = utf8_text
        del utf8_text
        del im
        gc.collect()
    # Final text in single file with a mark
    txt = ""
    mark = "\n### NEW PAGE ###\n"
    txt = mark.join(text[key] for key in text.keys()) + mark

    # delete temp folder
    shutil.rmtree(directory, ignore_errors=True)

    logging.info(
        f"Extract Image type: HTML, number pages: {len(text.keys())}, Length text: {len(txt[: -len(mark)])}"
    )
    del files
    del im_b64
    del images_arr
    del soup
    del html_data
    del text
    gc.collect()
    return txt[: -len(mark)]


def extract_text_from_odt(data, app):
    """
    extract and Process images in a odf file with images embedded in base64 in src and extract text via ocr
    params:
    data: list of documents from detect AI
    app global variables from FastAPI. contains tessetact and leptonica objects
    return:
    text extracted and endpoint
    """

    # page number
    n = 0
    # temporal directories
    directory = os.path.join(ROOT_DIR, "tmp")
    if not os.path.exists(directory):
        os.makedirs(directory)

    images = os.path.join(ROOT_DIR, "tmp", "images")
    if not os.path.exists(images):
        os.makedirs(images)

    text = {}
    im_b64 = data.get("source").get("content")
    # write docx to temporal folder
    odf_temp = os.path.join(directory, "test.odt")
    # write temp odf file
    with open(odf_temp, "wb") as f:
        f.write(base64.b64decode(im_b64))

    logging.info(f"Extract Image type: ODT, Length: {len(base64.b64decode(im_b64))}")
    # read the temporal odf file with zipfile
    with zipfile.ZipFile(odf_temp) as zf:
        # Get list of file names inside zip
        file_list = zf.namelist()
        # Check for Pictures folder
        pictures_folder = "Pictures"
        for file in file_list:
            if pictures_folder in file:
                # Extract Pictures folder
                zf.extract(file, path=images)
        del file_list
        gc.collect()

    # list all images extracted in images. Zip file look for images in Folder named Pictures and extract them under that name.
    # Odf it is a kind of zip container
    files = os.listdir(os.path.join(images, "Pictures"))

    for i, file in zip(range(len(files)), files):
        im = Image.open(os.path.join(images, "Pictures", file))
        # Extract text from Image
        utf8_text = ocr_image(im, app)
        text[i] = utf8_text
        del im
        del utf8_text
        gc.collect()
    # Final text in single file with a mark
    txt = ""
    mark = "\n### NEW PAGE ###\n"
    txt = mark.join(text[key] for key in text.keys()) + mark

    # delete temp folder
    shutil.rmtree(directory, ignore_errors=True)
    logging.info(
        f"Extract Image type: ODT, number pages: {len(text.keys())}, Length text: {len(txt[: -len(mark)])}"
    )
    del files
    del im_b64
    del text
    gc.collect()
    return txt[: -len(mark)]


def extract_text_from_rtf(data, app):
    """
    extract and Process images in a rtf file with images. Use of unrtf tool which has to be installed in the OS
    params:
    data: list of documents from detect AI
    app global variables from FastAPI. contains tessetact and leptonica objects
    return:
    text extracted and endpoint
    """

    def convert_rtf_to_doc(input_file, output_file, images, ROOT_DIR):
        os.chdir(images)
        os.system(f"/usr/bin/unrtf --html {input_file} > {output_file}")
        os.chdir(ROOT_DIR)

    # temporal directories
    directory = os.path.join(ROOT_DIR, "tmp")
    if not os.path.exists(directory):
        os.makedirs(directory)
    images = os.path.join(ROOT_DIR, "tmp", "images")
    if not os.path.exists(images):
        os.makedirs(images)

    text = {}
    im_b64 = data.get("source").get("content")
    # write docx to temporal folder
    rtf_temp = os.path.join(directory, "temp.rtf")

    with open(rtf_temp, "wb") as f:
        f.write(base64.b64decode(im_b64))

    logging.info(f"Extract Image type: RTF, Length: {len(base64.b64decode(im_b64))}")
    # convert RTF to HTML
    output_file = os.path.join(images, "test_rtf.html")
    convert_rtf_to_doc(rtf_temp, output_file, images, ROOT_DIR)
    # Open html and extract Images
    with open(output_file) as html_wr:
        html_data = html_wr.read()

    soup = BS(html_data, features="lxml")
    # find all references to Images  and save Pathnames to list files
    images_arr = soup.find_all("img")
    files = []
    for ind, i in zip(images_arr, range(len(images_arr))):
        files.append(ind["src"])

    # OCR each image on the list
    for i, file in zip(range(len(files)), files):
        im = Image.open(os.path.join(images, file))
        # Extract text from Image
        utf8_text = ocr_image(im, app)
        text[i] = utf8_text
        del im
        del utf8_text
        gc.collect()
    # Final text in single file with a mark
    txt = ""
    mark = "\n### NEW PAGE ###\n"
    txt = mark.join(text[key] for key in text.keys()) + mark

    # delete temp folder
    shutil.rmtree(directory, ignore_errors=True)
    logging.info(
        f"Extract Image type: RTF, number pages: {len(text.keys())}, Length text: {len(txt[: -len(mark)])}"
    )
    del files
    del im_b64
    del images_arr
    del soup
    del text
    gc.collect()
    return txt[: -len(mark)]


def extract_documents_from_request(document, app):
    """
    Process documents depending of their file type
    params:
    document: dictionary containing metadata and data of a document
    app global variables from FastAPI. contains tessetact and leptonica objects
    return:
    text extracted and endpoint

    """
    # get extension of document from request
    ext = document.get("source").get("file_type")

    if ext in ["jpg", "png", "jpeg", "gif", "bmp"]:
        utf8_text, document = extract_from_single_image(document, app)
    elif ext in ["tif", "tiff"]:
        utf8_text = extract_text_from_tiff(document, app)
    elif ext == "pdf":
        utf8_text = extract_text_from_pdf(document, app)
    elif ext in ["doc", "docx"]:
        utf8_text = extract_text_from_doc(document, app)
    elif ext in ["ppt", "pptx"]:
        utf8_text = extract_text_from_ppt(document, app)
    elif ext in ["htm", "html"]:
        utf8_text = extract_text_from_html(document, app)
    elif ext in ["odt"]:
        utf8_text = extract_text_from_odt(document, app)
    elif ext in ["rtf"]:
        utf8_text = extract_text_from_rtf(document, app)

    return utf8_text, document


def process_request(
    list_docs: list[dict],
    app: FastAPI,
    jobs: dict,
    new_task: Job,
    cypher: int = 0,
):
    """ "
    process list of base64 Image documents to OCR


    params:
    list_docs: list of documents from detect AI. List of dictionaries containing metadata and data of documents
    app global variables from FastAPI. contains tessetact and leptonica objects
    jobs: dictionary to hold status of the Job
    new_task: Job object
    cypher: boolean , True encrypt the output/False No
    use_rapid_ocr: string , YES use rapid ocr/NO No
    return:
    processed_docs : list of dictionaries containing document processed , this is pass trought OCR and text extracted. The extracted text replace the original base64 content
    documents_non_teathred : list of dictionaries containing {id : reason of not treating this id}
    """
    jobs[new_task.uid].status = "Start Extracting Text From Documents"
    documents_non_teathred = []
    processed_docs = []
    for data, i in zip(list_docs, range(len(list_docs))):
        file_type = data.get("source").get("file_type")
        file_name = data.get("source").get("file_name")
        file_uri = data.get("source").get("fs").get("uri")
        try:

            logging.info(
                f"Processing file : {file_name} length : {len(data.get('source').get('content'))}"
            )
            logging.info(f"Processing file : {file_name} uri : {file_uri}")
            # if file type not valid
            if not file_type or not (file_type in image_file_names):
                documents_non_teathred.append({data.get("id"): "File type not valid"})
                logging.warning(f"File : {file_name}  No Valid Extension")
                continue

            utf8_text, data = extract_documents_from_request(data, app)
            # if no text extracted doc to documents_non_teathred
            if len(utf8_text) < 3:
                documents_non_teathred.append(
                    {data.get("id"): "No Images in this document to process"}
                )
                logging.warning(f"File : {file_name}  No Text in this file")
                continue
            # create encrypted content

            if cypher:
                # chunck = encrypt(utf8_text)
                # data["crypto"] = (base64.b64encode(chunck)).decode("ascii")
                # case that we want to activate recover the file crypto.py , import it and
                logging.warning("this functionality has been deactivated")
            else:
                pass
            data["source"]["content"] = utf8_text
            # bug 6179
            data["source"]["fs"]["uri"] = data["source"]["fs"]["uri"].lower()
            # add to processed docs
            processed_docs.append(data)
            logging.info(f"Extracted Text From {len(processed_docs)} Documents")
            jobs[new_task.uid].status = (
                f"Extracted Text From {len(processed_docs)} Documents"
            )
        except Exception as e:
            documents_non_teathred.append(
                {
                    data.get(
                        "id"
                    ): f"No Images in this document to process and produced Exception {str(e)}"
                }
            )
            logging.error(f"File : {file_name}  produced Exception {str(e)}")
            if len(list_docs) > 1 and i + 1 < len(list_docs):
                continue

    return processed_docs, documents_non_teathred
