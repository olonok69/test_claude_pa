import base64
import io
import os
from PIL import Image, ImageSequence, ImageFile
from pdf2image import convert_from_bytes
import numpy as np
import gc
import shutil
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from bs4 import BeautifulSoup as BS
import zipfile
import logging
from rapidocr_onnxruntime import RapidOCR
import subprocess


ImageFile.LOAD_TRUNCATED_IMAGES = True

ROOT_DIR = os.path.dirname(os.path.abspath(__file__))


def create_engine_rapidocr(config_path: str = None):
    """
    Create RapidOCR Engine
    return:
    RapidOCR Engine
    """
    # create RapidOCR Engine
    engine = RapidOCR(config_path=config_path)
    return engine


def process_image_nats(image_data, pre_processing, engine):
    """Create an image from the binary data
    and apply any pre-processing if needed.
    Then pass the image to the OCR function.
    Args:
        image_data: binary data of the image
        pre_processing: function for pre-processing the image data (optional)
        engine: OCR engine or extraction engine for processing
    Returns:
        ocr_result: result of the OCR process
    """
    image = Image.open(io.BytesIO(image_data))

    # Apply any pre-processing if needed
    if "2x-zoom" in pre_processing:
        # Implement zoom functionality using PIL
        new_size = (image.width * 2, image.height * 2)
        image = image.resize(new_size, Image.LANCZOS)

    # Convert back to bytes if your ocr_image function accepts bytes
    with io.BytesIO() as output:
        image.save(output, format=image.format or "PNG")
        processed_image_data = output.getvalue()

    # Now pass this to your OCR function
    ocr_result = ocr_image_nats(processed_image_data, engine)
    return ocr_result


def ocr_to_text(ocr_output):
    """
    # Sort the OCR output based on the y-coordinate of the bounding box
    each element of the list is another list wihich contain 3 elements

        element 0 bounding box coordinates

        element 1 Text extracted

        element 2 score

        Example:

        [[[529.0, 149.0], [577.0, 149.0], [577.0, 174.0], [529.0, 174.0]],
        'R&D',
        0.9749892751375834]
    Args:
        ocr_output : list of lists, output from Rapid OCR
    """
    # Sort the OCR output based on the y-coordinate of the bounding box
    sorted_ocr = sorted(ocr_output, key=lambda x: x[0][0][1])

    # Initialize variables
    lines = []
    current_line = []
    current_y = sorted_ocr[0][0][0][1]

    # Iterate through the sorted OCR output
    for item in sorted_ocr:
        bbox, text, score = item
        bbox_y = bbox[0][1]

        # Check if the bounding box is on the same line as the current line
        if abs(bbox_y - current_y) <= 3:
            current_line.append(text + " ")
        else:
            lines.append(" ".join(current_line))
            current_line = [text + " "]
            current_y = bbox_y

    # Append the last line
    if current_line:
        lines.append(" ".join(current_line))

    # Join all lines into a single text output
    text_output = "\n".join(lines)
    return text_output


def ocr_image_nats(data, engine):
    """
    Extract text from an Image using Leptonica and tessetact
    params:
    im: PIL image
    app: FastAPI endpoint. Contains global objects with tesseract and Leptonica C objects
    return:
    text extracted
    """

    result, _ = engine(data)
    # Extract text while preserving spaces and returns
    extracted_text = ocr_to_text(result)  # "\n".join([line[1] for line in result])

    return extracted_text


def extract_from_single_image_data(file_data, engine, pre_processing):
    """
    Process imagenes and extract text via ocr
    params:
    file_data: binary data of the image
    engine: OCR engine or extraction engine for processing
    pre_processing: function for pre-processing the image data (optional)
    return:
    text extracted and endpoint
    """
    # Process the image data using the specified engine and pre-processing function
    extracted_text = process_image_nats(file_data, pre_processing, engine)
    return extracted_text


def extract_text_from_pdf_data(file_data, engine, pre_processing=[]):
    """
    Extract and process images in a PDF file and extract text via OCR using process_image_nats.
    Args:
        message_data: binary data from a NATS message
        pre_processing: list of pre-processing steps to apply
        engine: OCR engine for processing
    Returns:
        Extracted text with page separators
    """

    try:
        # Decode the message data
        img_bytes = file_data

        # Convert PDF to images
        pdf_file = convert_from_bytes(img_bytes, dpi=300)
        text = {}

        # Process each page in the PDF
        for i, page in enumerate(pdf_file):
            # Convert page to a binary image data
            page_arr = np.asarray(page)
            im = Image.fromarray(page_arr)

            # Convert PIL Image to binary data
            with io.BytesIO() as output:
                im.save(output, format="PNG")
                image_data = output.getvalue()

            # Extract text using process_image_nats
            utf8_text = process_image_nats(image_data, pre_processing, engine)
            text[i] = utf8_text

            # Clean up
            del im
            del page_arr
            del utf8_text
            gc.collect()

        # Combine all page texts with page separators
        mark = "\n\n"
        txt = mark.join(text[key] for key in text.keys()) + mark

        logging.info(
            f"Extract Image type: PDF, number pages: {len(text.keys())}, Length text: {len(txt[: -len(mark)])}"
        )

        # Clean up resources
        del pdf_file
        del img_bytes
        del text
        gc.collect()

        # Return the text without the trailing page marker
        return txt[: -len(mark)]

    except Exception as e:
        logging.error(f"Error extracting text from PDF: {str(e)}")
        return ""


def extract_text_from_tiff_data(file_data, engine, pre_processing=[]):
    """Extract text from a TIFF file using process_image_nats

    Args:
        file_data: binary data of the TIFF file
        engine: OCR engine for processing
        pre_processing: list of pre-processing steps to apply

    Returns:
        text: extracted text from all pages with page separators
    """
    try:
        # Open Image
        imfile = Image.open(io.BytesIO(file_data))
        text = {}

        # Process each page in the TIFF
        for i, page in enumerate(ImageSequence.Iterator(imfile)):
            # Convert to RGB mode (required for some processing)
            im = page.convert("RGB")

            # Convert PIL Image to binary data for process_image_nats
            with io.BytesIO() as output:
                im.save(output, format="PNG")
                image_data = output.getvalue()

            # Extract text using process_image_nats
            utf8_text = process_image_nats(image_data, pre_processing, engine)
            text[i] = utf8_text

            # Clean up to avoid memory issues
            del im
            del utf8_text
            gc.collect()

        # Combine all page texts with page separators
        mark = "\n\n"
        txt = mark.join(text[key] for key in text.keys()) + mark

        logging.info(
            f"Extract Image type: TIFF/TIF, number pages: {len(text.keys())}, Length text: {len(txt[: -len(mark)])}"
        )

        # Clean up resources
        del imfile
        del text
        gc.collect()

        # Return text without the trailing page marker
        return txt[: -len(mark)]

    except Exception as e:
        logging.error(f"Error extracting text from TIFF: {str(e)}")
        return ""


def extract_images_from_docx_data(docx_data, output_folder):
    """
    Extract images from a DOCX document using binary data

    Args:
        docx_data: Binary data of the .docx file
        output_folder: Path to the output folder where the images will be extracted

    Returns:
        List of paths to the extracted images
    """
    # Ensure the output folder exists
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    extracted_images = []

    try:
        # Create a BytesIO object from the binary data
        docx_buffer = io.BytesIO(docx_data)

        # Open the DOCX file as a zip file from the buffer
        with zipfile.ZipFile(docx_buffer, "r") as docx_zip:
            # Iterate through the files in the zip
            for file in docx_zip.namelist():
                # Check if the file is an image
                if file.startswith("word/media/"):
                    # Extract the image to the output folder
                    image_path = os.path.join(output_folder, os.path.basename(file))

                    # Extract the file content
                    image_data = docx_zip.read(file)

                    # Save the image to the output folder
                    with open(image_path, "wb") as f:
                        f.write(image_data)

                    logging.info(f"Extracted {file} to {image_path}")
                    extracted_images.append(image_path)

        return extracted_images

    except Exception as e:
        logging.error(f"Error extracting images from DOCX: {str(e)}")
        return []


def convert_doc_to_docx(doc_path, docx_path, directory):
    """
    # Convert DOC to DOCX using LibreOffice command line
    Args:
        doc_path: Path to the input .doc file.
        docx_path: Path to the output .docx file.
    """
    command = [
        "libreoffice",
        "--headless",
        "--convert-to",
        "docx",
        "--outdir",
        os.path.dirname(docx_path),
        doc_path,
    ]
    process = subprocess.Popen(command)
    try:
        process.wait(timeout=20)
        logging.info(
            f"Converted {doc_path} to {docx_path} susccessfully. convert_doc_to_docx Function"
        )
        process.kill()
    except subprocess.TimeoutExpired:
        # Kill the process if it exceeds the timeout
        process.kill()
        shutil.rmtree(directory, ignore_errors=True)
        logging.error(f"Error converting {doc_path} to {docx_path}: Timeout")
    except Exception as e:
        logging.error(f"Error converting {doc_path} to {docx_path}: {str(e)}")
        shutil.rmtree(directory, ignore_errors=True)


def extract_text_from_doc_data(file_data, engine, pre_processing=[]):
    """
    Extract and process images in a doc/docx file and extract text via OCR

    Args:
        file_data: Binary data of the doc/docx file
        engine: OCR engine for processing
        pre_processing: List of pre-processing steps to apply

    Returns:
        Extracted text from all images in the document with page separators
    """
    try:
        import tempfile

        # Create temporary directories
        directory = tempfile.mkdtemp()
        images_dir = os.path.join(directory, "images")
        os.makedirs(images_dir, exist_ok=True)

        # Determine file type by checking the binary signature
        if file_data[:8] == b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1":  # DOC file signature
            # Handle DOC file
            temp_doc = os.path.join(directory, "temp.doc")
            with open(temp_doc, "wb") as f:
                f.write(file_data)

            # Convert DOC to DOCX
            temp_docx = os.path.join(directory, "temp.docx")
            convert_doc_to_docx(temp_doc, temp_docx, directory)

            if not os.path.exists(temp_docx):
                raise Exception("Failed to convert DOC to DOCX")

        else:  # Assume DOCX
            # Write DOCX to temporal folder
            temp_docx = os.path.join(directory, "temp.docx")
            with open(temp_docx, "wb") as f:
                f.write(file_data)

        logging.info(f"Extract Image type: DOC/DOCX, Length: {len(file_data)}")

        # Extract images from DOCX
        extracted_images = extract_images_from_docx_data(file_data, images_dir)

        # Process each extracted image
        text = {}
        for i, image_path in enumerate(extracted_images):
            logging.info(f"Extracting text from {image_path}")

            # Read the image file
            with open(image_path, "rb") as img_file:
                image_data = img_file.read()

            # Process the image using process_image_nats
            utf8_text = process_image_nats(image_data, pre_processing, engine)
            text[i] = utf8_text

            # Clean up
            del utf8_text
            gc.collect()

        # If we have no images but have a text-based document, try to extract text directly
        if not text:
            try:
                import docx

                doc = docx.Document(temp_docx)
                plain_text = "\n".join(
                    [para.text for para in doc.paragraphs if para.text]
                )
                if plain_text:
                    text[0] = plain_text
            except Exception as e:
                logging.warning(f"Failed to extract text directly from DOCX: {str(e)}")

        # Final text in single file with a mark
        if text:
            mark = "\n\n"
            txt = mark.join(text[key] for key in text.keys()) + mark
            result = txt[: -len(mark)]

            logging.info(
                f"Extract Image type: DOC/DOCX, number pages: {len(text.keys())}, Length text: {len(result)}"
            )
        else:
            result = ""
            logging.warning("No text extracted from DOC/DOCX file")

        # Clean up temporary directory
        shutil.rmtree(directory, ignore_errors=True)

        # Clean up memory
        del text
        gc.collect()

        return result

    except Exception as e:
        logging.error(f"Error extracting text from DOC/DOCX: {str(e)}")
        # Clean up temporary files if they exist
        try:
            if "directory" in locals():
                shutil.rmtree(directory, ignore_errors=True)
        except:
            pass
        return ""


def convert_ppt_to_pptx_with_soffice(ppt_path, output_dir):
    """
    Converts a .ppt file to .pptx format using the soffice command-line tool.

    Args:
      ppt_path: Path to the input .ppt file.
      output_dir: Path to the output directory.

    Returns:
      True if the conversion was successful, False otherwise.
    """
    command = [
        "soffice",
        "--headless",
        "--convert-to",
        "pptx",
        "--outdir",
        output_dir,
        ppt_path,
    ]
    process = subprocess.Popen(command)
    try:
        process.wait(timeout=20)
        logging.info(
            f"Converted {ppt_path} to .pptx in {output_dir}. convert_ppt_to_pptx_with_soffice Function"
        )
        process.kill()

    except subprocess.TimeoutExpired:
        # Kill the process if it exceeds the timeout
        process.kill()
        logging.error(f"Error converting {ppt_path} to .pptx: Timeout")
    return


def extract_text_from_ppt_data(file_data, engine, pre_processing=[]):
    """
    Extract and process images in a ppt/pptx file and extract text via OCR using process_image_nats

    Args:
        file_data: binary data of the PPT/PPTX file
        engine: OCR engine for processing
        pre_processing: list of pre-processing steps to apply

    Returns:
        text extracted from all slides with page separators
    """

    def write_image(shape, images, n):
        image = shape.image
        # ---get image "file" contents---
        image_bytes = image.blob
        # ---make up a name for the file, e.g. 'image.jpg'---
        image_filename = os.path.join(images, "ppt_image{:03d}.{}".format(n, image.ext))
        n += 1

        with open(image_filename, "wb") as f:
            f.write(image_bytes)
        return n

    def visitor(shape, images, n):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for s in shape.shapes:
                n = visitor(s, images, n)
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            n = write_image(shape, images, n)
        return n

    def iter_picture_shapes(prs, n):
        for slide in prs.slides:
            for shape in slide.shapes:
                n = visitor(shape, images, n)

    try:
        # page number
        n = 0
        # temporal directories
        directory = os.path.join(ROOT_DIR, "tmp")
        if not os.path.exists(directory):
            os.makedirs(directory)

        images = os.path.join(ROOT_DIR, "tmp", "images")
        if not os.path.exists(images):
            os.makedirs(images)

        text = {}

        # Determine file extension based on byte signature or other metadata
        # For simplicity, we'll look at the beginning bytes to guess the format
        if file_data[:8] == b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1":  # PPT signature
            ext = "ppt"
        else:  # Default to PPTX
            ext = "pptx"

        # Process based on file type
        if ext == "pptx":
            # write pptx to temporal folder
            temp_pptx = os.path.join(directory, "temp.pptx")
            with open(temp_pptx, "wb") as f:
                f.write(file_data)
        elif ext == "ppt":
            # write ppt to temporal folder
            temp_ppt = os.path.join(directory, "temp.ppt")
            with open(temp_ppt, "wb") as f:
                f.write(file_data)
            temp_pptx = os.path.join(directory, "temp.pptx")
            # convert PPT to PPTX
            convert_ppt_to_pptx_with_soffice(temp_ppt, directory)

        logging.info(f"Extract Image type: PPT/PPTX, Length: {len(file_data)}")

        # Open PPT and extract Images
        iter_picture_shapes(Presentation(temp_pptx), n)

        # list all images extracted in images
        files = os.listdir(images)

        for i, file in zip(range(len(files)), files):
            # Read the image file into binary data
            with open(os.path.join(images, file), "rb") as img_file:
                image_data = img_file.read()

            # Process the image using process_image_nats instead of ocr_image
            utf8_text = process_image_nats(image_data, pre_processing, engine)
            text[i] = utf8_text
            del utf8_text
            gc.collect()

        # Final text in single file with a mark
        mark = "\n\n"
        txt = mark.join(text[key] for key in text.keys()) + mark

        # delete temp folder
        shutil.rmtree(directory, ignore_errors=True)

        logging.info(
            f"Extract Image type: PPT/PPTX, number pages: {len(text.keys())}, Length text: {len(txt[: -len(mark)])}"
        )
        del files
        del text
        gc.collect()
        return txt[: -len(mark)]

    except Exception as e:
        logging.error(f"Error extracting text from PPT/PPTX: {str(e)}")
        return ""


def extract_text_from_html_data(file_data, engine, pre_processing=[]):
    """
    Extract and process images embedded in an HTML file (base64 in src) and extract text via OCR

    Args:
        file_data: binary data of the HTML file
        engine: OCR engine for processing
        pre_processing: list of pre-processing steps to apply

    Returns:
        text extracted from all images with page separators
    """
    try:
        # page number
        n = 0
        # temporal directories
        directory = os.path.join(ROOT_DIR, "tmp")
        if not os.path.exists(directory):
            os.makedirs(directory)

        images = os.path.join(ROOT_DIR, "tmp", "images")
        if not os.path.exists(images):
            os.makedirs(images)

        text = {}

        # write html to temporal folder
        html_temp = os.path.join(directory, "temp.html")

        with open(html_temp, "wb") as f:
            f.write(file_data)

        logging.info(f"Extract Image type: HTML, Length: {len(file_data)}")

        # Open html and extract Images
        with open(html_temp, "r", encoding="utf-8") as html_wr:
            html_data = html_wr.read()

        soup = BS(html_data, features="lxml")
        images_arr = soup.find_all("img")

        for ind, i in zip(images_arr, range(len(images_arr))):
            # Skip images without src or base64 data
            if not ind.get("src") or "," not in ind["src"]:
                continue

            try:
                image_data_base64 = ind["src"].split(",")[1]
                decoded_img_data = base64.b64decode(image_data_base64)

                # Save the image temporarily
                image_path = os.path.join(images, f"site_{i}.png")
                with open(image_path, "wb+") as img_wr:
                    img_wr.write(decoded_img_data)

                # Process the image using process_image_nats instead of ocr_image
                utf8_text = process_image_nats(decoded_img_data, pre_processing, engine)
                text[i] = utf8_text

                del utf8_text
                del decoded_img_data
                gc.collect()
            except Exception as e:
                logging.error(f"Error processing image {i} in HTML: {str(e)}")
                continue

        logging.info(f"Extract Image type: HTML, number images: {len(images_arr)}")

        # Final text in single file with a mark
        mark = "\n\n"
        txt = mark.join(text[key] for key in text.keys()) + mark

        # delete temp folder
        shutil.rmtree(directory, ignore_errors=True)

        logging.info(
            f"Extract Image type: HTML, number pages: {len(text.keys())}, Length text: {len(txt[: -len(mark)])}"
        )

        # Clean up

        del images_arr
        del soup
        del html_data
        del text
        gc.collect()

        return txt[: -len(mark)]

    except Exception as e:
        logging.error(f"Error extracting text from HTML: {str(e)}")
        return ""


def extract_text_from_odt_data(file_data, engine, pre_processing=[]):
    """
    Extract and process images in an ODT file and extract text via OCR

    Args:
        file_data: binary data of the ODT file
        engine: OCR engine for processing
        pre_processing: list of pre-processing steps to apply

    Returns:
        text extracted from all images with page separators
    """
    try:
        # page number
        n = 0
        # temporal directories
        directory = os.path.join(ROOT_DIR, "tmp")
        if not os.path.exists(directory):
            os.makedirs(directory)

        images = os.path.join(ROOT_DIR, "tmp", "images")
        if not os.path.exists(images):
            os.makedirs(images)

        text = {}

        # write ODT to temporal folder
        odf_temp = os.path.join(directory, "test.odt")

        # write temp odf file
        with open(odf_temp, "wb") as f:
            f.write(file_data)

        logging.info(f"Extract Image type: ODT, Length: {len(file_data)}")

        # read the temporal odf file with zipfile
        with zipfile.ZipFile(odf_temp) as zf:
            # Get list of file names inside zip
            file_list = zf.namelist()
            # Check for Pictures folder
            pictures_folder = "Pictures"

            for file in file_list:
                if pictures_folder in file:
                    # Extract Pictures folder
                    zf.extract(file, path=images)

            del file_list
            gc.collect()

        # List all images extracted in images
        # ODT is a zip container with images in the Pictures folder
        pictures_path = os.path.join(images, "Pictures")

        if os.path.exists(pictures_path):
            files = os.listdir(pictures_path)

            for i, file in zip(range(len(files)), files):
                image_path = os.path.join(pictures_path, file)

                # Read the image file into binary data
                with open(image_path, "rb") as img_file:
                    image_data = img_file.read()

                # Process the image using process_image_nats instead of ocr_image
                utf8_text = process_image_nats(image_data, pre_processing, engine)
                text[i] = utf8_text

                del utf8_text
                gc.collect()

            # Final text in single file with a mark
            mark = "\n\n"
            txt = mark.join(text[key] for key in text.keys()) + mark

            # delete temp folder
            shutil.rmtree(directory, ignore_errors=True)

            logging.info(
                f"Extract Image type: ODT, number pages: {len(text.keys())}, Length text: {len(txt[: -len(mark)])}"
            )

            del files
            del text
            gc.collect()

            return txt[: -len(mark)]
        else:
            logging.warning("No Pictures folder found in the ODT file")
            return ""

    except Exception as e:
        logging.error(f"Error extracting text from ODT: {str(e)}")
        return ""


def extract_text_from_rtf_data(file_data, engine, pre_processing=[]):
    """
    Extract and process images in a RTF file and extract text via OCR

    Args:
        file_data: binary data of the RTF file
        engine: OCR engine for processing
        pre_processing: list of pre-processing steps to apply

    Returns:
        text extracted from all images with page separators
    """

    def convert_rtf_to_doc(input_file, output_file, images, ROOT_DIR):
        """Convert RTF to HTML using unrtf tool"""
        os.chdir(images)
        os.system(f"/usr/bin/unrtf --html {input_file} > {output_file}")
        os.chdir(ROOT_DIR)

    try:
        # temporal directories
        directory = os.path.join(ROOT_DIR, "tmp")
        if not os.path.exists(directory):
            os.makedirs(directory)

        images = os.path.join(ROOT_DIR, "tmp", "images")
        if not os.path.exists(images):
            os.makedirs(images)

        text = {}

        # write RTF to temporal folder
        rtf_temp = os.path.join(directory, "temp.rtf")
        with open(rtf_temp, "wb") as f:
            f.write(file_data)

        logging.info(f"Extract Image type: RTF, Length: {len(file_data)}")

        # convert RTF to HTML
        output_file = os.path.join(images, "test_rtf.html")
        convert_rtf_to_doc(rtf_temp, output_file, images, ROOT_DIR)

        # Open html and extract Images
        with open(output_file, "r", encoding="utf-8", errors="ignore") as html_wr:
            html_data = html_wr.read()

        soup = BS(html_data, features="lxml")

        # find all references to Images and save Pathnames to list files
        images_arr = soup.find_all("img")
        files = []
        for ind in images_arr:
            if ind.get("src"):
                files.append(ind["src"])

        # OCR each image on the list
        for i, file in zip(range(len(files)), files):
            try:
                # Read the image file into binary data
                image_path = os.path.join(images, file)
                with open(image_path, "rb") as img_file:
                    image_data = img_file.read()

                # Process the image using process_image_nats
                utf8_text = process_image_nats(image_data, pre_processing, engine)
                text[i] = utf8_text

                del utf8_text
                gc.collect()
            except Exception as img_error:
                logging.error(f"Error processing image {file} in RTF: {str(img_error)}")
                continue

        # Final text in single file with a mark
        mark = "\n\n"
        txt = mark.join(text[key] for key in text.keys()) + mark if text else ""

        # delete temp folder
        shutil.rmtree(directory, ignore_errors=True)

        logging.info(
            f"Extract Image type: RTF, number pages: {len(text.keys())}, Length text: {len(txt[: -len(mark) if text else 0])}"
        )

        del files
        del images_arr
        del soup
        del text
        gc.collect()

        return txt[: -len(mark)] if text else ""

    except Exception as e:
        logging.error(f"Error extracting text from RTF: {str(e)}")
        return ""


def extract_documents_from_file(file_data, mime_type, engine, pre_processing=None):
    """
    Process documents depending on their MIME type

    params:
    file_data: binary data of the document
    mime_type: MIME type of the document
    engine: OCR engine or extraction engine for processing
    pre_processing: function for pre-processing the image data (optional)

    return:
    extracted_text: text extracted from the document
    """
    # Map MIME types to file extensions
    mime_to_ext = {
        "image/jpeg": "jpg",
        "image/png": "png",
        "image/gif": "gif",
        "image/bmp": "bmp",
        "image/tiff": "tiff",
        "application/pdf": "pdf",
        "application/msword": "doc",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
        "application/vnd.ms-powerpoint": "ppt",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
        "text/html": "html",
        "application/vnd.oasis.opendocument.text": "odt",
        "application/rtf": "rtf",
        "text/rtf": "rtf",
    }

    # Get the file extension from the MIME type
    ext = mime_to_ext.get(mime_type, None)

    if not ext:
        logging.error(f"Unsupported MIME type: {mime_type}")
        return ""

    # Process the file based on its type
    if ext in ["jpg", "png", "jpeg", "gif", "bmp"]:
        extracted_text = extract_from_single_image_data(
            file_data, engine, pre_processing
        )
    elif ext in ["tif", "tiff"]:
        extracted_text = extract_text_from_tiff_data(file_data, engine, pre_processing)
    elif ext == "pdf":
        extracted_text = extract_text_from_pdf_data(file_data, engine, pre_processing)
    elif ext in ["doc", "docx"]:
        extracted_text = extract_text_from_doc_data(file_data, engine, pre_processing)
    elif ext in ["ppt", "pptx"]:
        extracted_text = extract_text_from_ppt_data(file_data, engine, pre_processing)
    elif ext in ["htm", "html"]:
        extracted_text = extract_text_from_html_data(file_data, engine, pre_processing)
    elif ext in ["odt"]:
        extracted_text = extract_text_from_odt_data(file_data, engine, pre_processing)
    elif ext in ["rtf"]:
        extracted_text = extract_text_from_rtf_data(file_data, engine, pre_processing)
    else:
        logging.error(f"Unsupported file extension: {ext}")
        extracted_text = ""

    return extracted_text


def determine_mime_type(file_data):
    """
    Determine the MIME type of a file from its binary data

    params:
    file_data: binary data of the file

    return:
    mime_type: detected MIME type
    """
    import magic

    mime = magic.Magic(mime=True)
    mime_type = mime.from_buffer(file_data)

    return mime_type


def process_request_nat(file_data, engine, pre_processing=None):
    """
    Process a document file and extract text from it based on its MIME type

    params:
    file_data: binary data of the document file
    engine: OCR engine or extraction engine
    pre_processing: optional pre-processing function to apply to the file data

    return:
    extracted_text: text extracted from the document
    """
    # Determine the MIME type of the file
    mime_type = determine_mime_type(file_data)
    logging.info(f"Detected MIME type: {mime_type}")

    # Extract text from the document based on its MIME type
    extracted_text = extract_documents_from_file(
        file_data, mime_type, engine, pre_processing
    )

    return extracted_text
